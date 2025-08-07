# ingestion_agent.py

import os
import pandas as pd
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

class IngestionAgent:
    SUPPORTED_TYPES = ('.pdf', '.docx', '.pptx', '.csv', '.txt', '.md')

    def __init__(self):
        pass

    def extract_text(self, file_path: str) -> str:
        ext = os.path.splitext(file_path)[-1].lower()

        if ext == ".pdf":
            reader = PdfReader(file_path)
            return "\n".join([page.extract_text() or "" for page in reader.pages])

        elif ext == ".docx":
            doc = Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])

        elif ext == ".pptx":
            prs = Presentation(file_path)
            return "\n".join([
                shape.text for slide in prs.slides
                for shape in slide.shapes if hasattr(shape, "text")
            ])

        elif ext == ".csv":
            df = pd.read_csv(file_path)
            return df.to_string(index=False)

        elif ext in [".txt", ".md"]:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()

        else:
            raise ValueError(f"Unsupported file format: {ext}")

    def chunk_text(self, text: str, chunk_size: int = 300, overlap: int = 50) -> list:
        words = text.split()
        chunks = []
        for i in range(0, len(words), chunk_size - overlap):
            chunk = " ".join(words[i:i + chunk_size])
            chunks.append(chunk)
        return chunks

    def ingest(self, file_paths: list[str]) -> dict:
        document_chunks = {}
        for path in file_paths:
            try:
                text = self.extract_text(path)
                chunks = self.chunk_text(text)
                document_chunks[path] = chunks
            except Exception as e:
                print(f"[Error] Failed to process {path}: {e}")
        return document_chunks